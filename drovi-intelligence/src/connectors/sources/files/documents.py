"""
Document Connector

Extracts text content from various document formats:
- PDF
- DOCX (Word)
- XLSX (Excel)
- PPTX (PowerPoint)
- TXT, MD, CSV

Supports both local files and URLs.
"""

import csv
import io
import os
import re
from collections.abc import AsyncIterator
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

import httpx
import structlog

from src.connectors.base.config import ConnectorConfig, StreamConfig, SyncMode
from src.connectors.base.connector import BaseConnector, RecordBatch, ConnectorRegistry
from src.connectors.base.records import RecordType
from src.connectors.base.state import ConnectorState
from src.connectors.http import request_with_retry

logger = structlog.get_logger()


@dataclass
class Document:
    """Represents an extracted document."""

    id: str
    name: str
    file_type: str
    source_path: str
    content: str
    page_count: int | None = None
    word_count: int = 0
    metadata: dict[str, Any] = field(default_factory=dict)
    created_at: str | None = None
    modified_at: str | None = None


class DocumentConnector(BaseConnector):
    """
    Connector for extracting content from documents.

    Supports:
    - PDF (via pypdf)
    - DOCX (via python-docx)
    - XLSX (via openpyxl)
    - PPTX (via python-pptx)
    - Plain text files (TXT, MD, CSV)

    Can read from:
    - Local file paths
    - URLs
    - S3-style paths (with appropriate credentials)
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "doc",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".txt": "txt",
        ".md": "markdown",
        ".csv": "csv",
        ".json": "json",
    }

    def __init__(self):
        """Initialize document connector."""
        pass

    async def check_connection(
        self,
        config: ConnectorConfig,
    ) -> tuple[bool, str | None]:
        """Check if document sources are accessible."""
        try:
            source_paths = config.get_setting("source_paths", [])
            source_urls = config.get_setting("source_urls", [])

            if not source_paths and not source_urls:
                return False, "No source_paths or source_urls configured"

            # Check local paths
            for path in source_paths:
                if not os.path.exists(path):
                    return False, f"Path does not exist: {path}"

            # Check URLs are reachable
            for url in source_urls[:1]:  # Just check first URL
                async with httpx.AsyncClient() as client:
                    response = await request_with_retry(
                        client,
                        "HEAD",
                        url,
                        follow_redirects=True,
                    )
                    if response.status_code >= 400:
                        return False, f"URL not accessible: {url}"

            return True, None

        except Exception as e:
            return False, f"Connection check failed: {str(e)}"

    async def discover_streams(
        self,
        config: ConnectorConfig,
    ) -> list[StreamConfig]:
        """Discover available document streams."""
        return [
            StreamConfig(
                stream_name="documents",
                sync_mode=SyncMode.FULL_REFRESH,
            ),
        ]

    async def read_stream(
        self,
        config: ConnectorConfig,
        stream: StreamConfig,
        state: ConnectorState,
    ) -> AsyncIterator[RecordBatch]:
        """Read documents from configured sources."""
        if stream.stream_name == "documents":
            async for batch in self._read_documents(config, stream, state):
                yield batch
        else:
            logger.warning(f"Unknown stream: {stream.stream_name}")

    async def _read_documents(
        self,
        config: ConnectorConfig,
        stream: StreamConfig,
        state: ConnectorState,
    ) -> AsyncIterator[RecordBatch]:
        """Read and extract content from documents."""
        source_paths = config.get_setting("source_paths", [])
        source_urls = config.get_setting("source_urls", [])
        recursive = config.get_setting("recursive", True)

        batch = self.create_batch(stream.stream_name, config.connection_id)

        # Process local paths
        for source_path in source_paths:
            path = Path(source_path)
            if path.is_file():
                doc = await self._process_file(path)
                if doc:
                    record = self.create_record(
                        record_id=doc.id,
                        stream_name=stream.stream_name,
                        data=self._document_to_record(doc),
                    )
                    record.record_type = RecordType.DOCUMENT
                    batch.add_record(record)
            elif path.is_dir():
                for file_path in self._iterate_files(path, recursive):
                    doc = await self._process_file(file_path)
                    if doc:
                        record = self.create_record(
                            record_id=doc.id,
                            stream_name=stream.stream_name,
                            data=self._document_to_record(doc),
                        )
                        record.record_type = RecordType.DOCUMENT
                        batch.add_record(record)

                    # Batch every 50 documents
                    if len(batch.records) >= 50:
                        batch.complete(has_more=True)
                        yield batch
                        batch = self.create_batch(stream.stream_name, config.connection_id)

        # Process URLs
        for url in source_urls:
            doc = await self._process_url(url)
            if doc:
                record = self.create_record(
                    record_id=doc.id,
                    stream_name=stream.stream_name,
                    data=self._document_to_record(doc),
                )
                record.record_type = RecordType.DOCUMENT
                batch.add_record(record)

        if batch.records:
            batch.complete(has_more=False)
            yield batch

    def _iterate_files(
        self,
        directory: Path,
        recursive: bool,
    ) -> list[Path]:
        """Iterate over files in a directory."""
        files = []
        pattern = "**/*" if recursive else "*"

        for file_path in directory.glob(pattern):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                files.append(file_path)

        return sorted(files)

    async def _process_file(self, file_path: Path) -> Document | None:
        """Process a single file and extract content."""
        try:
            extension = file_path.suffix.lower()
            file_type = self.SUPPORTED_EXTENSIONS.get(extension)

            if not file_type:
                logger.debug(f"Unsupported file type: {extension}")
                return None

            # Get file metadata
            stat = file_path.stat()
            created_at = datetime.fromtimestamp(stat.st_ctime).isoformat()
            modified_at = datetime.fromtimestamp(stat.st_mtime).isoformat()

            # Extract content based on type
            content, page_count, metadata = await self._extract_content(
                file_path, file_type
            )

            if not content:
                return None

            return Document(
                id=str(file_path.absolute()),
                name=file_path.name,
                file_type=file_type,
                source_path=str(file_path.absolute()),
                content=content,
                page_count=page_count,
                word_count=len(content.split()),
                metadata=metadata,
                created_at=created_at,
                modified_at=modified_at,
            )

        except Exception as e:
            logger.error(f"Failed to process file: {file_path}", error=str(e))
            return None

    async def _process_url(self, url: str) -> Document | None:
        """Download and process a document from URL."""
        try:
            async with httpx.AsyncClient() as client:
                response = await request_with_retry(
                    client,
                    "GET",
                    url,
                    follow_redirects=True,
                )
                response.raise_for_status()

                # Determine file type from URL or content-type
                content_type = response.headers.get("content-type", "")
                file_name = url.split("/")[-1].split("?")[0]
                extension = Path(file_name).suffix.lower()

                file_type = self.SUPPORTED_EXTENSIONS.get(extension)
                if not file_type:
                    # Try to guess from content-type
                    if "pdf" in content_type:
                        file_type = "pdf"
                    elif "word" in content_type or "docx" in content_type:
                        file_type = "docx"
                    elif "text/plain" in content_type:
                        file_type = "txt"
                    elif "text/csv" in content_type:
                        file_type = "csv"
                    else:
                        logger.debug(f"Unknown content type: {content_type}")
                        return None

                # Extract content from bytes
                content, page_count, metadata = await self._extract_content_from_bytes(
                    response.content, file_type
                )

                if not content:
                    return None

                return Document(
                    id=url,
                    name=file_name,
                    file_type=file_type,
                    source_path=url,
                    content=content,
                    page_count=page_count,
                    word_count=len(content.split()),
                    metadata=metadata,
                )

        except Exception as e:
            logger.error(f"Failed to process URL: {url}", error=str(e))
            return None

    async def _extract_content(
        self,
        file_path: Path,
        file_type: str,
    ) -> tuple[str, int | None, dict[str, Any]]:
        """Extract content from a file."""
        with open(file_path, "rb") as f:
            content_bytes = f.read()

        return await self._extract_content_from_bytes(content_bytes, file_type)

    async def _extract_content_from_bytes(
        self,
        content_bytes: bytes,
        file_type: str,
    ) -> tuple[str, int | None, dict[str, Any]]:
        """Extract content from bytes based on file type."""
        metadata: dict[str, Any] = {}
        page_count = None
        content = ""

        if file_type == "pdf":
            content, page_count, metadata = self._extract_pdf(content_bytes)
        elif file_type in ("docx", "doc"):
            content, metadata = self._extract_docx(content_bytes)
        elif file_type in ("xlsx", "xls"):
            content, metadata = self._extract_xlsx(content_bytes)
        elif file_type in ("pptx", "ppt"):
            content, page_count, metadata = self._extract_pptx(content_bytes)
        elif file_type in ("txt", "markdown"):
            content = content_bytes.decode("utf-8", errors="replace")
        elif file_type == "csv":
            content, metadata = self._extract_csv(content_bytes)
        elif file_type == "json":
            content = content_bytes.decode("utf-8", errors="replace")
            metadata["format"] = "json"

        return content, page_count, metadata

    def _extract_pdf(
        self,
        content_bytes: bytes,
    ) -> tuple[str, int, dict[str, Any]]:
        """Extract text from PDF."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(content_bytes))
            pages = []
            metadata = {}

            # Extract metadata
            if reader.metadata:
                metadata = {
                    "title": reader.metadata.get("/Title"),
                    "author": reader.metadata.get("/Author"),
                    "subject": reader.metadata.get("/Subject"),
                    "creator": reader.metadata.get("/Creator"),
                }

            # Extract text from each page
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)

            return "\n\n".join(pages), len(reader.pages), metadata

        except ImportError:
            logger.warning("pypdf not installed, cannot extract PDF")
            return "", 0, {}
        except Exception as e:
            logger.error(f"Failed to extract PDF: {e}")
            return "", 0, {}

    def _extract_docx(
        self,
        content_bytes: bytes,
    ) -> tuple[str, dict[str, Any]]:
        """Extract text from DOCX."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(io.BytesIO(content_bytes))
            paragraphs = []
            metadata = {}

            # Extract core properties
            if doc.core_properties:
                metadata = {
                    "title": doc.core_properties.title,
                    "author": doc.core_properties.author,
                    "subject": doc.core_properties.subject,
                    "keywords": doc.core_properties.keywords,
                }

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    paragraphs.append(row_text)

            return "\n\n".join(paragraphs), metadata

        except ImportError:
            logger.warning("python-docx not installed, cannot extract DOCX")
            return "", {}
        except Exception as e:
            logger.error(f"Failed to extract DOCX: {e}")
            return "", {}

    def _extract_xlsx(
        self,
        content_bytes: bytes,
    ) -> tuple[str, dict[str, Any]]:
        """Extract text from XLSX."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
            sheets_content = []
            metadata = {"sheet_names": wb.sheetnames}

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v for v in row_values):
                        rows.append(" | ".join(row_values))

                if rows:
                    sheets_content.append(f"## {sheet_name}\n\n" + "\n".join(rows))

            return "\n\n".join(sheets_content), metadata

        except ImportError:
            logger.warning("openpyxl not installed, cannot extract XLSX")
            return "", {}
        except Exception as e:
            logger.error(f"Failed to extract XLSX: {e}")
            return "", {}

    def _extract_pptx(
        self,
        content_bytes: bytes,
    ) -> tuple[str, int, dict[str, Any]]:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content_bytes))
            slides_content = []
            metadata = {}

            # Extract core properties
            if prs.core_properties:
                metadata = {
                    "title": prs.core_properties.title,
                    "author": prs.core_properties.author,
                    "subject": prs.core_properties.subject,
                }

            # Extract text from slides
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)

                if slide_texts:
                    slides_content.append(f"## Slide {i}\n\n" + "\n".join(slide_texts))

            return "\n\n".join(slides_content), len(prs.slides), metadata

        except ImportError:
            logger.warning("python-pptx not installed, cannot extract PPTX")
            return "", 0, {}
        except Exception as e:
            logger.error(f"Failed to extract PPTX: {e}")
            return "", 0, {}

    def _extract_csv(
        self,
        content_bytes: bytes,
    ) -> tuple[str, dict[str, Any]]:
        """Extract text from CSV."""
        try:
            text = content_bytes.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text))
            rows = list(reader)

            metadata = {
                "row_count": len(rows),
                "column_count": len(rows[0]) if rows else 0,
            }

            # Format as table
            content = "\n".join(" | ".join(row) for row in rows)

            return content, metadata

        except Exception as e:
            logger.error(f"Failed to extract CSV: {e}")
            return "", {}

    def _document_to_record(self, doc: Document) -> dict[str, Any]:
        """Convert Document to a record dict."""
        # Extract first 500 chars as preview
        preview = doc.content[:500] + "..." if len(doc.content) > 500 else doc.content

        return {
            "id": doc.id,
            "type": f"document_{doc.file_type}",
            "name": doc.name,
            "file_type": doc.file_type,
            "source_path": doc.source_path,
            "content": doc.content,
            "preview": preview,
            "page_count": doc.page_count,
            "word_count": doc.word_count,
            "char_count": len(doc.content),
            "metadata": doc.metadata,
            "title": doc.metadata.get("title"),
            "author": doc.metadata.get("author"),
            "created_at": doc.created_at,
            "modified_at": doc.modified_at,
            "extracted_at": datetime.utcnow().isoformat(),
        }


# Register connector
ConnectorRegistry.register("documents", DocumentConnector)
